#!/usr/bin/env python3
"""EgoPed ONE-PAGER: the whole project on a single 16:9 page (paper + new baselines).

Numbers are the verified unified-evaluation values from research/findings.md.
Usage:  python research/src/build_onepager.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "research/to_human/EgoPed_onepager.pptx"
NAVY = RGBColor(0x11, 0x2B, 0x4A); BLUE = RGBColor(0x1E, 0x5E, 0xA8)
ACC = RGBColor(0x2E, 0x86, 0xDE); INK = RGBColor(0x1F, 0x2A, 0x37)
MUT = RGBColor(0x5B, 0x66, 0x72); LIGHT = RGBColor(0xEF, 0xF4, 0xFB)
W = RGBColor(0xFF, 0xFF, 0xFF); GOOD = RGBColor(0x15, 0x7A, 0x38)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW = prs.slide_width
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = W


def rect(x, y, w, h, c, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False; return sh


def tb(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = s.shapes.add_textbox(x, y, w, h).text_frame
    t.word_wrap = True; t.vertical_anchor = anchor
    t.margin_left = Pt(2); t.margin_right = Pt(2); t.margin_top = Pt(1); t.margin_bottom = Pt(1)
    return t


def run(p, txt, size, color=INK, bold=False, italic=False):
    r = p.add_run(); r.text = txt; r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = "Calibri"; return r


def sec(x, y, w, title):
    """Section heading with accent rule."""
    t = tb(x, y, w, Inches(0.26))
    run(t.paragraphs[0], title.upper(), 10.5, BLUE, bold=True)
    rect(x, y + Inches(0.245), w, Pt(1.4), ACC)


def bul(x, y, w, h, items, size=10):
    t = tb(x, y, w, h)
    for i, it in enumerate(items):
        p = t.paragraphs[0] if i == 0 else t.add_paragraph()
        p.space_after = Pt(3.5)
        run(p, "▪  " + it, size, INK)


def table(x, y, w, headers, rows, colw, fs=9.5, good=(), rh=0.235):
    g = s.shapes.add_table(len(rows) + 1, len(headers), x, y, w, Inches(rh) * (len(rows) + 1))
    t = g.table
    tot = sum(colw)
    for j, cw in enumerate(colw):
        t.columns[j].width = int(w * cw / tot)
    for j, h in enumerate(headers):
        c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.margin_left = Pt(4); c.margin_top = Pt(0); c.margin_bottom = Pt(0)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        run(p, h, fs, W, bold=True)
    for i, rv in enumerate(rows):
        for j, v in enumerate(rv):
            c = t.cell(i + 1, j); c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xDC, 0xEE, 0xE1) if i in good else (LIGHT if i % 2 == 0 else W)
            c.margin_left = Pt(4); c.margin_top = Pt(0); c.margin_bottom = Pt(0)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run(p, v, fs, GOOD if i in good else INK, bold=(i in good))
    return g


# ══ header ═════════════════════════════════════════════════════════════════
rect(0, 0, SW, Inches(0.95), NAVY)
rect(0, Inches(0.95), SW, Pt(3.5), ACC)
t = tb(Inches(0.35), Inches(0.04), Inches(8.85), Inches(0.88), MSO_ANCHOR.MIDDLE)
run(t.paragraphs[0], "EgoPed", 25, W, bold=True)
run(t.paragraphs[0], "   ·   Ego-Conditioned Pedestrian Motion Generation via Latent Diffusion", 12.5,
    RGBColor(0xCF, 0xE1, 0xF5))
p = t.add_paragraph()
run(p, "3D full-body motion conditioned on vehicle ego trajectory  ·  MLD backbone  ·  "
       "AVA + nuScenes + Waymo (11,954)", 9.5, RGBColor(0x9E, 0xC4, 0xEC))
# headline chip
rect(Inches(9.25), Inches(0.17), Inches(3.75), Inches(0.62), RGBColor(0x1B, 0x4E, 0x2E))
t = tb(Inches(9.3), Inches(0.19), Inches(3.65), Inches(0.58), MSO_ANCHOR.MIDDLE)
p = t.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "FID 2.880", 19, RGBColor(0x8F, 0xE3, 0xAC), bold=True)
run(p, "   −40% vs pooled", 12, RGBColor(0xC9, 0xF2, 0xD8))

LX, LW = Inches(0.35), Inches(6.45)
RX, RW = Inches(7.05), Inches(5.93)

# ══ left: key idea ═════════════════════════════════════════════════════════
sec(LX, Inches(1.12), LW, "The idea")
bul(LX, Inches(1.45), LW, Inches(0.95), [
    "Replace MLD's text conditioning with the vehicle ego trajectory.",
    "Feed the denoiser ALL 196 per-timestep ego tokens — not one pooled summary — so every "
    "denoising step has temporal resolution over the trajectory.",
    "Ego encoder is contrastively pretrained, then FROZEN; only the denoiser trains.",
], 10)

# ══ left: main results ═════════════════════════════════════════════════════
sec(LX, Inches(2.52), LW, "Main results  ·  unified evaluation, 3 replications")
table(LX, Inches(2.85), LW,
      ["System", "FID ↓", "R@1", "MM", "Div"],
      [["EgoPed-IA  (+ interaction sampling)", "2.880", "0.323", "3.37", "5.67"],
       ["EgoPed  (full-sequence conditioning)", "3.392", "0.548", "3.96", "5.79"],
       ["Unfrozen encoder", "5.079", "0.352", "2.62", "6.06"],
       ["Pooled-token baseline", "5.620", "0.665", "3.71", "5.82"],
       ["Larger latent (8 vs 4)", "6.677", "0.663", "3.37", "5.89"],
       ["Ground truth", "—", "—", "—", "5.50"]],
      [4.3, 1.5, 1.3, 1.2, 1.2], 9.5, good=(0,))

# ══ left: architecture-vs-data 2x2 ═════════════════════════════════════════
sec(LX, Inches(4.72), LW, "Architecture vs. data curation  ·  2×2")
table(LX, Inches(5.05), LW,
      ["FID ↓", "− interaction sampling", "+ interaction sampling"],
      [["Pooled", "6.214", "5.620"],
       ["Full-sequence", "3.392", "2.880"]],
      [2.6, 3.0, 3.0], 9.5, good=(1,))
t = tb(LX, Inches(5.85), LW, Inches(0.5))
run(t.paragraphs[0], "Conditioning granularity dominates (≈2.7 FID in both columns, ~5× the data effect); "
    "the two effects are additive → the headline is conservative, not inflated.", 9.5, MUT, italic=True)

# ══ right: ablations ═══════════════════════════════════════════════════════
sec(RX, Inches(1.12), RW, "What actually matters")
table(RX, Inches(1.45), RW,
      ["Question", "Finding", "Verdict"],
      [["Unfreeze the ego encoder?", "MM 3.96→2.62", "freezing critical"],
       ["Skip encoder pretraining?", "R@1 = chance", "≈ unconditional"],
       ["Cross- instead of self-attention?", "3.49 vs 3.34", "no difference (p=.33)"],
       ["Bigger VAE latent (8 vs 4)?", "6.68 vs 5.62", "19% worse"],
       ["Different training seed?", "3.39/4.72/4.80", "real variance"]],
      [3.6, 2.3, 2.9], 9.5)

# ══ right: baseline ladder ═════════════════════════════════════════════════
sec(RX, Inches(3.05), RW, "External baselines  ·  capability ladder")
table(RX, Inches(3.38), RW,
      ["Approach", "FID ↓", "Reading"],
      [["EgoPed-IA (ours)", "2.88", "in-domain, ego-conditioned"],
       ["Unconditional prior", "5.18", "conditioning worth 34%"],
       ["Pretrained text-to-motion MLD", "14.07", "T2M does not transfer"],
       ["Raw-space diffusion (MDM-style)", "26.43", "latent backbone justified"],
       ["Ego→motion regressor", "35.47", "L2-optimal, unrealistic"],
       ["Trajectory + static body", "47.7", "needs articulated motion"]],
      [3.6, 1.3, 3.5], 9.5, good=(0,))

# ══ right: validation axes ═════════════════════════════════════════════════
sec(RX, Inches(5.15), RW, "Validated on four independent axes")
bul(RX, Inches(5.48), RW, Inches(1.3), [
    "FID — 40% over pooled, non-overlapping CIs, unified eval pipeline.",
    "Held-out split — scene-disjoint, never used for selection: ranking preserved.",
    "Trajectory ADE/FDE — no learned evaluator: beats pooled on all 4 metrics "
    "despite lower R@1 → the R@1 gap is an embedding artifact.",
    "Behavioral probe — best coupling of stop/walk decisions to the ego; "
    "unconditional model is ego-blind (sanity ✓).",
], 9.5)

# ══ footer ═════════════════════════════════════════════════════════════════
rect(0, Inches(6.85), SW, Inches(0.65), LIGHT)
rect(0, Inches(6.85), SW, Pt(1.5), ACC)
t = tb(Inches(0.35), Inches(6.88), Inches(12.6), Inches(0.6), MSO_ANCHOR.MIDDLE)
p = t.paragraphs[0]
run(p, "IN FLIGHT: ", 9.5, BLUE, bold=True)
run(p, "information ladder — how much conditioning information is actually needed? "
       "0 bits (unconditional) → 2.99 bits (body text) → 6.09 bits (ego text) → continuous 196×2 trajectory (ours). ",
    9.5, INK)
run(p, "Open: ", 9.5, BLUE, bold=True)
run(p, "pose-label quality audit · EgoPed-IA is single-seed · no physics constraints · no perceptual study yet.",
    9.5, INK)

os.makedirs("research/to_human", exist_ok=True)
prs.save(OUT)
print("SAVED", OUT)
